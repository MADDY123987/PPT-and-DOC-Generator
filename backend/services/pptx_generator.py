from pptx import Presentation
from pptx.util import Inches, Pt
import os
import requests
from io import BytesIO
from urllib.parse import urlparse
import re  # for cleaning URLs

# Folder where ppt1.pptx ... ppt5.pptx live
TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), "ppt_templates")

# Map theme IDs to actual template files
TEMPLATE_MAP = {
  "ppt1": os.path.join(TEMPLATE_DIR, "ppt1.pptx"),
  "ppt2": os.path.join(TEMPLATE_DIR, "ppt2.pptx"),
  "ppt3": os.path.join(TEMPLATE_DIR, "ppt3.pptx"),
  "ppt4": os.path.join(TEMPLATE_DIR, "ppt4.pptx"),
  "ppt5": os.path.join(TEMPLATE_DIR, "ppt5.pptx"),
  "ppt6": os.path.join(TEMPLATE_DIR, "ppt6.pptx"),
  "ppt7": os.path.join(TEMPLATE_DIR, "ppt7.pptx"),
  "ppt8": os.path.join(TEMPLATE_DIR, "ppt8.pptx"),
  "ppt9": os.path.join(TEMPLATE_DIR, "ppt9.pptx"),
}



def _get_layout(prs: Presentation, index: int, fallback: int = 0):
    """Safely get a slide layout by index with fallback."""
    try:
        return prs.slide_layouts[index]
    except IndexError:
        return prs.slide_layouts[fallback]


def _remove_all_slides(prs: Presentation):
    """Remove all existing slides from a Presentation (keep theme)."""
    slide_ids = list(prs.slides._sldIdLst)  # internal list of slide IDs
    for slide_id in slide_ids:
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(slide_id)


def _split_into_paragraphs(text: str, max_sentences_per_para: int = 3):
    """Split a long caption into smaller paragraphs (by sentence) for better layout."""
    if not text:
        return []
    parts = [p.strip() for p in text.replace("\n", " ").split(".") if p.strip()]
    paragraphs = []
    buf = []
    for i, p in enumerate(parts):
        buf.append(p + ".")
        if (i + 1) % max_sentences_per_para == 0:
            paragraphs.append(" ".join(buf).strip())
            buf = []
    if buf:
        paragraphs.append(" ".join(buf).strip())
    return paragraphs


def _get_tmp_image_path(img_url: str, presentation_id: int, slide_index: int) -> str:
    """
    Download image from a URL (or use local path) and return a temp file path.
    Raises if download fails.
    """
    os.makedirs("storage", exist_ok=True)

    # Local path
    if not img_url.startswith("http"):
        if not os.path.isfile(img_url):
            raise RuntimeError(f"Local image not found: {img_url}")
        return os.path.abspath(img_url)

    # Remote URL
    headers = {"User-Agent": "Mozilla/5.0 (compatible; PPTGenerator/1.0)"}
    resp = requests.get(img_url, headers=headers, timeout=15)
    resp.raise_for_status()

    parsed = urlparse(img_url)
    path = parsed.path or ""
    ext = os.path.splitext(path)[1] or ".jpg"

    tmp_path = os.path.abspath(f"./storage/tmp_img_{presentation_id}_{slide_index}{ext}")
    with open(tmp_path, "wb") as f:
        f.write(resp.content)

    return tmp_path


def build_pptx(presentation_id: int, slides: list, config: dict, **kwargs) -> str:
    """
    Build a PPTX using one of the PowerPoint templates in services/ppt_templates.

    slides: list of dicts like:
      { "layout": "title"|"bullet"|"two_column"|"image", ... }

    config: dict containing styling:
      { "theme_id": "ppt1" | ... "ppt5" | None, ... }
    """

    # 1) Choose template
    theme_id = (config or {}).get("theme_id") or "ppt1"
    template_path = TEMPLATE_MAP.get(theme_id)

    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        _remove_all_slides(prs)
    else:
        prs = Presentation()

    # 2) Build slides
    for slide_data in slides:
        layout_type = slide_data.get("layout", "title")
        title_text = slide_data.get("title", "")

        if layout_type == "title":
            layout = _get_layout(prs, 0)
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = title_text or "Title"

        elif layout_type == "bullet":
            layout = _get_layout(prs, 1, fallback=0)
            slide = prs.slides.add_slide(layout)

            if slide.shapes.title:
                slide.shapes.title.text = title_text or ""

            body_placeholder = None
            for shp in slide.placeholders:
                if (
                    getattr(shp, "is_placeholder", False)
                    and getattr(shp, "placeholder_format", None)
                    and shp.placeholder_format.type not in (1,)
                ):
                    body_placeholder = shp
                    break

            bullets = slide_data.get("bullets", []) or []

            if body_placeholder is not None:
                tf = body_placeholder.text_frame
                tf.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                        p.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet

        elif layout_type == "two_column":
            layout = _get_layout(prs, 3, fallback=1)
            slide = prs.slides.add_slide(layout)

            if slide.shapes.title:
                slide.shapes.title.text = title_text or ""

            left_text = slide_data.get("left", "")
            right_text = slide_data.get("right", "")

            content_placeholders = [
                shp
                for shp in slide.placeholders
                if getattr(shp, "is_placeholder", False)
                and getattr(shp, "placeholder_format", None)
                and shp.placeholder_format.type not in (1,)
            ]

            if len(content_placeholders) >= 1:
                left_tf = content_placeholders[0].text_frame
                left_tf.clear()
                left_tf.paragraphs[0].text = left_text

            if len(content_placeholders) >= 2:
                right_tf = content_placeholders[1].text_frame
                right_tf.clear()
                right_tf.paragraphs[0].text = right_text

        elif layout_type == "image":
            layout = _get_layout(prs, 3, fallback=1)  # two-content
            slide_index = len(prs.slides)
            slide = prs.slides.add_slide(layout)

            if slide.shapes.title:
                slide.shapes.title.text = title_text or ""

            content_placeholders = [
                shp
                for shp in slide.placeholders
                if getattr(shp, "is_placeholder", False)
                and getattr(shp, "placeholder_format", None)
                and shp.placeholder_format.type not in (1,)
            ]
            img_placeholder = content_placeholders[0] if len(content_placeholders) >= 1 else None
            text_placeholder = content_placeholders[1] if len(content_placeholders) >= 2 else None

            img_url = slide_data.get("image_url")
            caption = slide_data.get("caption") or slide_data.get("description") or ""

            if img_url:
                img_url = re.sub(r"\s+", "", str(img_url))
                img_url = img_url.strip("()[]{}.,;")

            text_to_use = caption or title_text or ""

            # IMAGE
            if img_url:
                try:
                    tmp_path = _get_tmp_image_path(img_url, presentation_id, slide_index)

                    if img_placeholder is not None:
                        left = img_placeholder.left
                        top = img_placeholder.top
                        width = img_placeholder.width
                        height = img_placeholder.height

                        slide.shapes.add_picture(tmp_path, left, top, width=width, height=height)
                        try:
                            img_placeholder.text = ""
                        except Exception:
                            pass
                    else:
                        left = int(prs.slide_width * 0.08)
                        top = int(prs.slide_height * 0.25)
                        width = int(prs.slide_width * 0.4)
                        slide.shapes.add_picture(tmp_path, left, top, width=width)

                except Exception as e:
                    text_to_use = f"{caption or title_text or ''}\n\n(Image failed to load: {img_url})"
                    print("Image download/insert failed:", e)

            # TEXT
            paragraphs = _split_into_paragraphs(text_to_use, max_sentences_per_para=2)

            if text_placeholder is not None:
                tf = text_placeholder.text_frame
                tf.clear()
                if paragraphs:
                    tf.paragraphs[0].text = paragraphs[0]
                    for para in paragraphs[1:]:
                        p = tf.add_paragraph()
                        p.text = para
                else:
                    tf.paragraphs[0].text = text_to_use
            else:
                left = int(prs.slide_width * 0.55)
                top = int(prs.slide_height * 0.25)
                width = int(prs.slide_width * 0.35)
                height = int(prs.slide_height * 0.5)
                caption_box = slide.shapes.add_textbox(left, top, width, height)
                tf = caption_box.text_frame
                tf.clear()
                if paragraphs:
                    tf.paragraphs[0].text = paragraphs[0]
                    for para in paragraphs[1:]:
                        p = tf.add_paragraph()
                        p.text = para
                else:
                    tf.paragraphs[0].text = text_to_use

        else:
            layout = _get_layout(prs, 0)
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = title_text or "Slide"

    # 3) Save
    os.makedirs("storage", exist_ok=True)
    path = os.path.abspath(f"./storage/presentation_{presentation_id}.pptx")
    prs.save(path)
    return path
